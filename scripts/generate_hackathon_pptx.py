"""Generate ForgeCTO hackathon pitch deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parents[1] / "ForgeCTO_Hackathon_Pitch.pptx"

# Palette — teal/ink (avoid generic purple AI look)
INK = RGBColor(0x0F, 0x1C, 0x2E)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0A, 0x6F, 0x66)
SLATE = RGBColor(0x3D, 0x4F, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF4, 0xF7, 0xF8)
ACCENT = RGBColor(0xE8, 0xA3, 0x17)
LIGHT_LINE = RGBColor(0xD0, 0xDB, 0xE4)


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=SLATE, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_footer(slide, page, total=12):
    add_text_box(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        "ForgeCTO  ·  Mini Hackathon",
        size=11,
        color=SLATE,
    )
    add_text_box(
        slide,
        Inches(8.5),
        Inches(7.05),
        Inches(1),
        Inches(0.3),
        f"{page}/{total}",
        size=11,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )


def title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), TEAL)
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.55), title, size=32, bold=True, color=INK)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4), subtitle, size=16, color=SLATE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 12

    # 1 — Title
    s = prs.slides.add_slide(blank)
    add_bg(s, INK)
    add_rect(s, Inches(0), Inches(0), Inches(0.22), Inches(7.5), TEAL)
    add_text_box(s, Inches(0.9), Inches(1.8), Inches(11), Inches(0.4), "MINI HACKATHON", size=14, bold=True, color=TEAL)
    add_text_box(s, Inches(0.9), Inches(2.3), Inches(11), Inches(1), "ForgeCTO", size=54, bold=True, color=WHITE, font="Calibri")
    add_text_box(
        s,
        Inches(0.9),
        Inches(3.4),
        Inches(11),
        Inches(0.8),
        "Your virtual startup CTO — from one idea to architecture,\nschema, APIs, AWS plan, and a build roadmap.",
        size=20,
        color=RGBColor(0xC5, 0xD0, 0xDA),
    )
    add_text_box(s, Inches(0.9), Inches(5.6), Inches(11), Inches(0.4), "AI Startup CTO Agent", size=16, color=TEAL)

    # 2 — Problem
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "The Problem", "Early founders need CTO-level clarity — before they can hire one")
    cards = [
        ("No technical co-founder", "Ideas stall because founders don’t know what to build first or how to structure the system."),
        ("Scattered AI advice", "ChatGPT gives fragments. No coherent architecture, schema, APIs, or sprint plan."),
        ("Expensive delay", "Wrong stack or vague MVP burns months — or consulting fees founders can’t afford yet."),
    ]
    for i, (h, body) in enumerate(cards):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(1.8), Inches(3.9), Inches(4.2), WHITE)
        add_rect(s, x, Inches(1.8), Inches(3.9), Inches(0.12), TEAL)
        add_text_box(s, x + Inches(0.25), Inches(2.2), Inches(3.4), Inches(1), h, size=20, bold=True, color=INK)
        add_text_box(s, x + Inches(0.25), Inches(3.3), Inches(3.4), Inches(2.2), body, size=15, color=SLATE)
    add_footer(s, 2, total)

    # 3 — Solution
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Our Solution: ForgeCTO", "One idea in → a complete CTO pack out")
    add_text_box(
        s,
        Inches(0.6),
        Inches(1.7),
        Inches(12),
        Inches(0.6),
        "A multi-agent virtual CTO that turns a startup idea into a structured, reviewable technical plan.",
        size=18,
        color=SLATE,
    )
    outputs = [
        "Market research & competitors",
        "Prioritized MVP features",
        "System architecture (+ diagrams)",
        "Database schema",
        "API endpoints",
        "AWS design & cost estimate",
        "Roadmap & sprint plan",
        "Markdown + GitHub export",
    ]
    for i, item in enumerate(outputs):
        col = i % 4
        row = i // 4
        x = Inches(0.6 + col * 3.15)
        y = Inches(2.6 + row * 1.5)
        add_rect(s, x, y, Inches(2.95), Inches(1.2), WHITE)
        add_rect(s, x, y, Inches(0.12), Inches(1.2), TEAL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.35), Inches(2.5), Inches(0.6), item, size=15, bold=True, color=INK)
    add_footer(s, 3, total)

    # 4 — How it works
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "How It Works", "7 specialized agents in a LangGraph pipeline")
    steps = [
        ("1", "Planner", "Scope product"),
        ("2", "Research", "Market & features"),
        ("3", "Architecture", "System design"),
        ("4", "Database", "Schema design"),
        ("5", "API", "Endpoints"),
        ("6", "AWS", "Cloud & cost"),
        ("7", "Docs", "Roadmap & sprints"),
    ]
    for i, (n, title, sub) in enumerate(steps):
        x = Inches(0.35 + i * 1.85)
        add_rect(s, x, Inches(2.2), Inches(1.7), Inches(2.8), WHITE)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), Inches(2.45), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = n
        set_run(run, size=18, bold=True, color=WHITE)
        add_text_box(s, x + Inches(0.08), Inches(3.4), Inches(1.55), Inches(0.5), title, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.08), Inches(3.9), Inches(1.55), Inches(0.7), sub, size=12, color=SLATE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(s, x + Inches(1.55), Inches(2.55), Inches(0.35), Inches(0.4), "→", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(
        s,
        Inches(0.6),
        Inches(5.4),
        Inches(12),
        Inches(0.8),
        "Live progress streams to the UI (SSE). Each step writes structured artifacts into a shared project state.",
        size=16,
        color=SLATE,
    )
    add_footer(s, 4, total)

    # 5 — Demo flow
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Live Demo Flow", "What judges will see in 2 minutes")
    demo_steps = [
        ("01", "Enter a startup idea", "e.g. “Uber for pet grooming”"),
        ("02", "Watch agents run live", "Planner → Research → Architecture → …"),
        ("03", "Browse the CTO pack", "Tabs: research, schema, APIs, AWS, roadmap"),
        ("04", "Export & ship", "Download markdown or create GitHub issues"),
    ]
    for i, (n, h, d) in enumerate(demo_steps):
        y = Inches(1.7 + i * 1.15)
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE)
        add_rect(s, Inches(0.6), y, Inches(1.1), Inches(1.0), TEAL)
        add_text_box(s, Inches(0.7), y + Inches(0.28), Inches(0.9), Inches(0.5), n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(2.0), y + Inches(0.15), Inches(10), Inches(0.4), h, size=20, bold=True, color=INK)
        add_text_box(s, Inches(2.0), y + Inches(0.55), Inches(10), Inches(0.35), d, size=15, color=SLATE)
    add_footer(s, 5, total)

    # 6 — Tech stack
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Tech Stack", "Built for real agents — not just a chat wrapper")
    stacks = [
        ("Backend", ["FastAPI", "LangGraph + LangChain", "Gemini / OpenAI", "Tavily search", "SQLAlchemy + SQLite/Postgres"]),
        ("Frontend", ["React + TypeScript", "Vite", "Mermaid.js diagrams", "SSE live progress", "Clerk auth"]),
        ("Ops & export", ["Docker Compose", "Rate limiting", "Demo fallback mode", "Markdown download", "GitHub issue export"]),
    ]
    for i, (title, items) in enumerate(stacks):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(1.7), Inches(3.9), Inches(4.5), WHITE)
        add_rect(s, x, Inches(1.7), Inches(3.9), Inches(0.55), TEAL)
        add_text_box(s, x + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.4), title, size=18, bold=True, color=WHITE)
        add_bullets(s, x + Inches(0.25), Inches(2.5), Inches(3.4), Inches(3.4), items, size=16, color=SLATE)
    add_footer(s, 6, total)

    # 7 — Key features
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Why This Matters", "Product strengths for founders")
    feats = [
        ("Structured, not chatty", "Pydantic schemas force consistent CTO artifacts every run."),
        ("Specialized agents", "Each node has one job — better focus than one mega-prompt."),
        ("Always demoable", "Seed project + offline demo pack if LLM quota fails."),
        ("From plan to backlog", "Export docs and push GitHub issues for the engineering team."),
        ("Transparent progress", "Founders see each CTO step complete in real time."),
        ("Human-in-the-loop ready", "Draft for review — not blind automation of hard decisions."),
    ]
    for i, (h, d) in enumerate(feats):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(1.7 + row * 2.3)
        add_rect(s, x, y, Inches(3.9), Inches(2.05), WHITE)
        add_rect(s, x, y, Inches(3.9), Inches(0.1), ACCENT if row == 0 else TEAL)
        add_text_box(s, x + Inches(0.25), y + Inches(0.35), Inches(3.4), Inches(0.5), h, size=17, bold=True, color=INK)
        add_text_box(s, x + Inches(0.25), y + Inches(0.95), Inches(3.4), Inches(0.9), d, size=14, color=SLATE)
    add_footer(s, 7, total)

    # 8 — Differentiation
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Differentiation", "We don’t replace builders — we brief them")
    rows = [
        ("ChatGPT / Cursor", "Great for code & answers", "No end-to-end CTO pack or pipeline"),
        ("v0 / Lovable", "Generate UI fast", "Skip research, schema, AWS, sprints"),
        ("Notion / Miro AI", "General assistants", "Not a specialist technical workflow"),
        ("ForgeCTO", "Idea → full CTO pack", "Research → architecture → APIs → cloud → roadmap"),
    ]
    # header
    add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.55), TEAL_DARK)
    for i, h in enumerate(["Tool", "Strength", "Gap / ForgeCTO edge"]):
        add_text_box(s, Inches(0.8 + i * 4.0), Inches(1.8), Inches(3.7), Inches(0.4), h, size=15, bold=True, color=WHITE)
    for r, (a, b, c) in enumerate(rows):
        y = Inches(2.35 + r * 0.95)
        bg = WHITE if r % 2 == 0 else RGBColor(0xE8, 0xF0, 0xF2)
        if r == 3:
            bg = RGBColor(0xD4, 0xF0, 0xEC)
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.9), bg)
        vals = [a, b, c]
        for i, v in enumerate(vals):
            add_text_box(
                s,
                Inches(0.8 + i * 4.0),
                y + Inches(0.25),
                Inches(3.7),
                Inches(0.5),
                v,
                size=14,
                bold=(r == 3 or i == 0),
                color=INK if r == 3 else SLATE,
            )
    add_footer(s, 8, total)

    # 9 — Architecture snapshot
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "System Snapshot", "Simple architecture judges can follow")
    boxes = [
        (0.6, "React UI", "Idea form\nLive steps\nArtifact tabs"),
        (3.7, "FastAPI", "Projects API\nSSE events\nExport routes"),
        (6.8, "LangGraph", "7 agent nodes\nShared state\nRetries / timeouts"),
        (9.9, "Storage", "Postgres / SQLite\nArtifacts JSON\nRun events"),
    ]
    for x, title, body in boxes:
        add_rect(s, Inches(x), Inches(2.0), Inches(2.8), Inches(3.2), WHITE)
        add_rect(s, Inches(x), Inches(2.0), Inches(2.8), Inches(0.65), INK)
        add_text_box(s, Inches(x + 0.15), Inches(2.12), Inches(2.5), Inches(0.45), title, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(x + 0.2), Inches(2.95), Inches(2.4), Inches(2.0), body, size=15, color=SLATE, align=PP_ALIGN.CENTER)
    for x in (3.2, 6.3, 9.4):
        add_text_box(s, Inches(x), Inches(3.3), Inches(0.5), Inches(0.4), "→", size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(
        s,
        Inches(0.6),
        Inches(5.5),
        Inches(12),
        Inches(0.7),
        "LLM: Gemini (preferred) or OpenAI  ·  Research: Tavily (optional)  ·  Auth: Clerk",
        size=15,
        color=SLATE,
    )
    add_footer(s, 9, total)

    # 10 — Business & impact
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Who Pays & Why", "Lightweight go-to-market for a hackathon build")
    left = [
        "Primary users: first-time / early founders",
        "Secondary: accelerators & student cohorts",
        "Model: freemium runs + paid credits",
        "B2B: white-label CTO packs for programs",
    ]
    right = [
        "Value: CTO draft in minutes, not weeks",
        "Complements (not replaces) advisors",
        "Makes human consulting 10× more focused",
        "Low per-run cost with Gemini Flash tier",
    ]
    add_rect(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(4.3), WHITE)
    add_rect(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(0.6), TEAL)
    add_text_box(s, Inches(0.85), Inches(1.9), Inches(5.4), Inches(0.4), "Audience & model", size=18, bold=True, color=WHITE)
    add_bullets(s, Inches(0.9), Inches(2.7), Inches(5.3), Inches(3.0), left, size=16)

    add_rect(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.3), WHITE)
    add_rect(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(0.6), INK)
    add_text_box(s, Inches(7.05), Inches(1.9), Inches(5.4), Inches(0.4), "Impact", size=18, bold=True, color=WHITE)
    add_bullets(s, Inches(7.1), Inches(2.7), Inches(5.3), Inches(3.0), right, size=16)
    add_footer(s, 10, total)

    # 11 — Roadmap
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    title_bar(s, "Built Now / Next", "Honest scope for judges")
    add_rect(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(4.5), WHITE)
    add_rect(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(0.6), TEAL)
    add_text_box(s, Inches(0.85), Inches(1.9), Inches(5.4), Inches(0.4), "Shipped for hackathon", size=18, bold=True, color=WHITE)
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.7),
        Inches(5.3),
        Inches(3.3),
        [
            "Full 7-agent LangGraph pipeline",
            "Live SSE dashboard + Mermaid views",
            "Structured artifacts & seed demo",
            "Markdown + GitHub export",
            "Gemini/OpenAI + demo fallback",
        ],
        size=16,
    )
    add_rect(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.5), WHITE)
    add_rect(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(0.6), ACCENT)
    add_text_box(s, Inches(7.05), Inches(1.9), Inches(5.4), Inches(0.4), "Next 90 days", size=18, bold=True, color=INK)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.7),
        Inches(5.3),
        Inches(3.3),
        [
            "Approve / regenerate single steps",
            "Quality gates between agents",
            "Stronger research grounding",
            "Auth-scoped multi-user projects",
            "Accelerator cohort dashboard",
        ],
        size=16,
    )
    add_footer(s, 11, total)

    # 12 — Close / Q&A
    s = prs.slides.add_slide(blank)
    add_bg(s, INK)
    add_rect(s, Inches(0), Inches(0), Inches(0.22), Inches(7.5), TEAL)
    add_text_box(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.8), "Thank you", size=44, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.9),
        Inches(3.0),
        Inches(11),
        Inches(1.0),
        "ForgeCTO — hire a virtual CTO before you hire a real one.",
        size=22,
        color=RGBColor(0xC5, 0xD0, 0xDA),
    )
    add_text_box(s, Inches(0.9), Inches(4.4), Inches(11), Inches(0.5), "Questions?", size=28, bold=True, color=TEAL)
    add_text_box(
        s,
        Inches(0.9),
        Inches(5.3),
        Inches(11),
        Inches(0.8),
        "Demo tip: open seed project first, then run a fresh idea live.",
        size=16,
        color=RGBColor(0x9A, 0xAB, 0xBA),
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
